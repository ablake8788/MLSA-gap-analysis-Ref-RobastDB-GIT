# pptx_renderer.py          # markdown_to_pptx_table_style

from __future__ import annotations

import re
from typing import List, Optional

from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

from tga_cli.domain.models import RunContext, RunArtifacts


def _is_md_table_row(line: str) -> bool:
    s = line.strip()
    return s.startswith("|") and s.endswith("|") and "|" in s[1:-1]


def _is_md_table_sep(line: str) -> bool:
    s = line.strip()
    if not (s.startswith("|") and s.endswith("|")):
        return False
    cells = [c.strip() for c in s.strip("|").split("|")]
    if not cells:
        return False
    for c in cells:
        c2 = c.replace(":", "").replace("-", "")
        if c2 != "":
            return False
    return True


def _split_md_row(line: str) -> list[str]:
    return [c.strip() for c in line.strip().strip("|").split("|")]


def _chunk_table_for_slides(table_data: List[List[str]], max_body_rows: int) -> List[List[List[str]]]:
    if not table_data:
        return []
    cols = max(len(r) for r in table_data)
    normalized = [(r + [""] * cols)[:cols] for r in table_data]
    header = normalized[0]
    body = normalized[1:]
    if not body:
        return [[header]]

    max_body_rows = max(1, int(max_body_rows))
    chunks: List[List[List[str]]] = []
    for start in range(0, len(body), max_body_rows):
        chunks.append([header] + body[start:start + max_body_rows])
    return chunks


def _create_table_slide(prs: Presentation, title: str, table_data: list[list[str]]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Title box (non-clipping)
    slide_w = prs.slide_width
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.2), slide_w - Inches(1.2), Inches(0.95))
    tf = title_box.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = (title or "Section").strip()
    p.font.bold = True
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.LEFT

    rows = len(table_data)
    if rows == 0:
        return
    cols = max(len(r) for r in table_data)
    data = [(r + [""] * cols)[:cols] for r in table_data]

    top = Inches(0.82)
    left = Inches(0.5)
    width = slide_w - Inches(1.0)
    height = int((prs.slide_height - top - Inches(0.6)) / 3.3)

    tbl_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    tbl = tbl_shape.table

    # Column widths: favor 3-col compare
    if cols == 3:
        tbl.columns[0].width = int(width * 0.22)
        tbl.columns[1].width = int(width * 0.39)
        tbl.columns[2].width = int(width * 0.39)
    else:
        for j in range(cols):
            tbl.columns[j].width = int(width / cols)

    for i in range(rows):
        for j in range(cols):
            cell = tbl.cell(i, j)
            ctf = cell.text_frame
            ctf.clear()
            ctf.word_wrap = True
            ctf.vertical_anchor = MSO_ANCHOR.TOP
            para = ctf.paragraphs[0]
            para.text = (data[i][j] or "").strip()
            para.alignment = PP_ALIGN.LEFT
            if i == 0:
                para.font.bold = True
                para.font.size = Pt(10)
            else:
                para.font.bold = False
                para.font.size = Pt(8)


class PptxRenderer:
    def __init__(self, *, max_table_body_rows_per_slide: int = 12):
        self.max_rows = max_table_body_rows_per_slide

    def render(self, *, report_md: str, ctx: RunContext, artifacts: RunArtifacts) -> None:
        prs = Presentation()

        # Title slide
        s0 = prs.slides.add_slide(prs.slide_layouts[0])
        s0.shapes.title.text = "Competitor Gap Analysis"
        s0.placeholders[1].text = f"Primary Competitor: {ctx.competitor_url} | Generated: {ctx.generated_at}"

        lines = report_md.splitlines()
        i = 0
        current_section_title: Optional[str] = None
        pending_lines: List[str] = []

        def flush_pending():
            nonlocal pending_lines
            if not pending_lines:
                return
            rows = [["Item", "Notes"]]
            for t in pending_lines:
                rows.append([t, ""])
            for idx, chunk in enumerate(_chunk_table_for_slides(rows, self.max_rows)):
                suffix = " (cont.)" if idx > 0 else ""
                _create_table_slide(prs, (current_section_title or "Section") + suffix, chunk)
            pending_lines = []

        while i < len(lines):
            raw = lines[i].rstrip("\n")
            s = raw.strip()
            ######
            # Skip markdown horizontal rules so they don't become slide items
            if s == "---":
                i += 1
                continue

            #####

            if s.startswith("# "):
                flush_pending()
                current_section_title = s[2:].strip()
                i += 1
                continue
            if s.startswith("## "):
                flush_pending()
                current_section_title = s[3:].strip()
                i += 1
                continue
            if s.startswith("### "):
                flush_pending()
                current_section_title = s[4:].strip()
                i += 1
                continue

            if _is_md_table_row(raw) and (i + 1) < len(lines) and _is_md_table_sep(lines[i + 1]):
                flush_pending()
                table_lines = [raw, lines[i + 1]]
                i += 2
                while i < len(lines) and _is_md_table_row(lines[i]):
                    table_lines.append(lines[i])
                    i += 1

                header = _split_md_row(table_lines[0])
                body = [_split_md_row(x) for x in table_lines[2:] if _is_md_table_row(x)]
                table_data = [header] + body

                chunks = _chunk_table_for_slides(table_data, self.max_rows)
                for idx, chunk in enumerate(chunks):
                    suffix = " (cont.)" if idx > 0 else ""
                    _create_table_slide(prs, (current_section_title or "Section") + suffix, chunk)
                continue

            if s.startswith(("- ", "* ")):
                pending_lines.append(s[2:].strip())
            else:
                m = re.match(r"^\d+\.\s+(.*)$", s)
                if m:
                    pending_lines.append(m.group(1).strip())
                elif s:
                    pending_lines.append(s)

            i += 1

        flush_pending()
        prs.save(str(artifacts.pptx_path))
